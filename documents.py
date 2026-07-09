# documents.py —— 文档理解引擎（多格式解析 + 结构化提取 + 跨文档检索/比对）
#
# 对应需求：
#   1) 几乎全格式上传：PDF / Word / Excel / PPT / CSV / HTML / 网页 URL / Markdown / 文本 / JSON
#      + 扫描件 / 图片走「可选 OCR」（本沙箱无 tesseract，代码路径保留，自动优雅降级）
#   2) 百万字级长文档：按「结构化分块」处理（按标题/段落切分 + 重叠），不吞内存、可检索
#   3) 精准提取：表格、条款（第X条/Article X/1.1）、层级结构（标题树）
#   4) 跨文档关联比对与交叉验证：对所有文档的分块做向量检索 + 相似配对 + 数值差异扫描
#
# 设计原则：
#   - 不依赖大模型也能跑：解析/提取/检索全部用解析器 + 启发式算法；
#     若 .env 配了真实嵌入/LLM，则自动升级为语义能力（复用 memory.py 的 embed/cosine）。
#   - 每个解析器用「懒加载 + try/except」，缺某个库不会让整个模块崩，只会降级并留 notes。

import os
import re
import json
import uuid
import datetime

# 复用记忆库的「余弦」作为跨文档检索底座（单一真相来源，避免重复造轮子）。
# 注意：文档检索的向量用「本地关键词向量」(见 _doc_vec)，刻意不依赖 LLM embedding 网络——
# 这样百万字级文档在离线/MOCK 环境下也能秒级解析与检索，不会因为逐个 chunk 调网络 embedding
# 而卡死（真实环境下若要更强语义，可把 _doc_vec 换成 memory.embed 的批量调用）。
from memory import _cosine
from storage import DATA_DIR, UPLOAD_DIR

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DOCS_FILE = os.path.join(DATA_DIR, "documents.json")


def _doc_vec(text):
    """文档检索用的本地关键词向量（中文按单字、英文/数字按词，小写化）。
    纯本地、零依赖、离线可用——保证大规模文档检索不依赖 LLM embedding 网络。
    若要升级为真实语义向量，把这里换成 memory.embed 的批量调用即可（接口一致）。"""
    text = (text or "").lower()
    d = {}
    for tk in re.findall(r"[a-z0-9]+|[一-鿿]", text):
        d[tk] = d.get(tk, 0) + 1
    return d

# 分块参数：单块上限 + 重叠（重叠让跨块句子不被切断，检索更准）
CHUNK_SIZE = 1400
CHUNK_OVERLAP = 200

# 支持的扩展名 -> 逻辑类型
EXT_MAP = {
    "pdf": "pdf", "doc": "docx", "docx": "docx",
    "xls": "xlsx", "xlsx": "xlsx",
    "ppt": "pptx", "pptx": "pptx",
    "csv": "csv", "tsv": "csv",
    "html": "html", "htm": "html",
    "md": "md", "markdown": "md",
    "txt": "txt", "text": "txt",
    "json": "json",
    "png": "image", "jpg": "image", "jpeg": "image",
    "gif": "image", "bmp": "image", "webp": "image",
    "tif": "image", "tiff": "image",
}

# ============================================================
# 存取（文档库 = documents.json）
# ============================================================

def _ensure_dirs():
    os.makedirs(UPLOAD_DIR, exist_ok=True)


def _load_docs():
    if not os.path.exists(DOCS_FILE):
        return {}
    try:
        with open(DOCS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _save_docs(docs):
    _ensure_dirs()
    with open(DOCS_FILE, "w", encoding="utf-8") as f:
        json.dump(docs, f, ensure_ascii=False, indent=2)


# ============================================================
# 解析器（懒加载，缺库就降级）
# ============================================================

def _parse_pdf(path):
    """PDF：用 pdfplumber 抽文字 + 表格；pypdf 作兜底。返回 (text, tables, meta)。"""
    text, tables, meta = "", [], {}
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            meta["pages"] = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                t = page.extract_text() or ""
                text += t + "\n"
                for tbl in (page.extract_tables() or []):
                    if tbl:
                        tables.append(_normalize_table(tbl, page=i))
    except Exception:
        # 兜底：pypdf 只抽文字，拿不到表格
        try:
            from pypdf import PdfReader
            r = PdfReader(path)
            meta["pages"] = len(r.pages)
            text = "\n".join((p.extract_text() or "") for p in r.pages)
        except Exception as e:
            meta["notes"] = f"PDF 解析失败：{e}"
    # 扫描件检测：文字太少则提示（OCR 见 _ocr_text）
    if len(text.strip()) < 50:
        ocr = _ocr_text(path, is_pdf=True)
        if ocr:
            text = ocr
            meta["ocr"] = True
        else:
            meta.setdefault("notes", "")
            meta["notes"] += "（疑似扫描件：本环境未安装 OCR，无法提取文字；在你本机装 tesseract 后可识别）"
    return text, tables, [], meta


def _parse_docx(path):
    """Word：保留标题层级（Heading 样式）+ 正文 + 表格。"""
    text, tables, structure, meta = [], [], [], {}
    try:
        from docx import Document
        doc = Document(path)
        for p in doc.paragraphs:
            style = (p.style.name or "") if p.style else ""
            line = p.text
            text.append(line)
            m = re.match(r"Heading\s*(\d+)", style)
            if m and line.strip():
                structure.append({"level": int(m.group(1)), "title": line.strip()})
            elif style.startswith("Title") and line.strip():
                structure.append({"level": 1, "title": line.strip()})
        for tbl in doc.tables:
            rows = [[c.text for c in r.cells] for r in tbl.rows]
            if rows:
                tables.append(_normalize_table(rows))
    except Exception as e:
        meta["notes"] = f"Word 解析失败：{e}"
    return "\n".join(text), tables, structure, meta


def _parse_xlsx(path):
    """Excel：每个 sheet 当一张表，正文为各 cell 拼接。"""
    import pandas as pd
    tables, meta = [], {}
    frames = pd.read_excel(path, sheet_name=None, engine="openpyxl")
    meta["sheets"] = list(frames.keys())
    parts = []
    for name, df in frames.items():
        df = df.fillna("")
        headers = [str(h) for h in df.columns]
        rows = df.astype(str).values.tolist()
        tables.append({"caption": f"Sheet: {name}", "headers": headers, "rows": rows})
        parts.append(f"[{name}]\n" + df.to_string(index=False))
    return "\n".join(parts), tables, [], meta


def _parse_pptx(path):
    """PPT：逐幻灯片抽文字 + 表格。"""
    text, tables, meta = [], [], {}
    try:
        from pptx import Presentation
        prs = Presentation(path)
        meta["slides"] = len(prs.slides)
        for i, slide in enumerate(prs.slides, 1):
            buf = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    buf.append(shape.text_frame.text)
                if shape.has_table:
                    tbl = shape.table
                    rows = [[c.text for c in r.cells] for r in tbl.rows]
                    if rows:
                        tables.append(_normalize_table(rows, page=i))
            text.append("\n".join(buf))
    except Exception as e:
        meta["notes"] = f"PPT 解析失败：{e}"
    return "\n".join(text), tables, [], meta


def _parse_csv(path):
    """CSV/TSV：整体作为一张表。"""
    import pandas as pd
    sep = "\t" if path.lower().endswith(".tsv") else ","
    df = pd.read_csv(path, sep=sep, dtype=str).fillna("")
    headers = [str(h) for h in df.columns]
    rows = df.astype(str).values.tolist()
    tables = [{"caption": os.path.basename(path), "headers": headers, "rows": rows}]
    return df.to_string(index=False), tables, [], {}


def _parse_html(path_or_html, is_path=True):
    """HTML / 网页：trafilatura 抽正文；表格用 bs4 兜底抽。"""
    html = None
    if is_path:
        with open(path_or_html, "r", encoding="utf-8", errors="ignore") as f:
            html = f.read()
    else:
        html = path_or_html
    text, tables = "", []
    try:
        import trafilatura
        text = trafilatura.extract(html, include_tables=False) or ""
    except Exception:
        text = ""
    # 表格用 bs4 兜底
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        for tbl in soup.find_all("table")[:20]:
            rows = []
            for tr in tbl.find_all("tr"):
                cells = [c.get_text(strip=True) for c in tr.find_all(["td", "th"])]
                if cells:
                    rows.append(cells)
            if rows:
                tables.append(_normalize_table(rows))
    except Exception:
        pass
    if not text:
        # trafilatura 失败就退回纯文本
        try:
            from bs4 import BeautifulSoup
            text = BeautifulSoup(html, "html.parser").get_text("\n", strip=True)
        except Exception:
            text = re.sub(r"<[^>]+>", "", html)
    return text, tables, [], {}


def _parse_plain(path, as_markdown=False):
    """txt / md / json：直接读。md 额外抽 # 标题做结构。"""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    structure = []
    if as_markdown:
        for line in text.splitlines():
            m = re.match(r"^(#{1,6})\s+(.*)", line)
            if m:
                structure.append({"level": len(m.group(1)), "title": m.group(2).strip()})
    if path.lower().endswith(".json"):
        try:
            obj = json.loads(text)
            text = json.dumps(obj, ensure_ascii=False, indent=2)
        except Exception:
            pass
    return text, [], structure, {}


# ============================================================
# OCR（可选；无 tesseract 时降级）
# ============================================================

def _ocr_available():
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def _ocr_text(path, is_pdf=False):
    """对图片或 PDF 首页做 OCR。无 tesseract 返回空字符串。"""
    if not _ocr_available():
        return ""
    try:
        import pytesseract
        from PIL import Image
        if is_pdf:
            try:
                from pdf2image import convert_from_path
                imgs = convert_from_path(path, first_page=1, last_page=1)
                if imgs:
                    return pytesseract.image_to_string(imgs[0], lang="chi_sim+eng")
            except Exception:
                return ""
        else:
            return pytesseract.image_to_string(Image.open(path), lang="chi_sim+eng")
    except Exception:
        return ""


def _parse_image(path):
    """图片：尝试 OCR；不可用则只记录元信息。"""
    text = _ocr_text(path, is_pdf=False)
    meta = {}
    if text:
        meta["ocr"] = True
    else:
        meta["notes"] = "本环境未安装 OCR（tesseract），无法识别图片中的文字；在你本机装 tesseract + pytesseract 后可识别。"
    return text, [], [], meta


# ============================================================
# 表格/结构/条款 的通用处理
# ============================================================

def _normalize_table(rows, page=None):
    """把不规则二维数组规整成 {headers, rows}，首行当表头（若像表头）。"""
    if not rows:
        return {"headers": [], "rows": []}
    # 补齐每行列数
    ncol = max(len(r) for r in rows)
    norm = [r + [""] * (ncol - len(r)) for r in rows]
    headers = norm[0]
    body = norm[1:]
    # 若首行与后续行格式无差异，宁可把首行也当数据（交给上层判断）
    return {"caption": f"第 {page} 页" if page else "", "headers": headers, "rows": body}


def _detect_headings(text):
    """格式无关标题检测：返回 [{level, title, char_start}]。启发式 + 正则。
    O(n) 实现：用预计算的「下一非空行」数组替代 lines.index（O(n^2)）。"""
    headings = []
    lines = text.splitlines()
    pos = 0
    n = len(lines)
    # 预计算每个位置之后第一个非空行（从后往前扫一遍）
    next_nonempty = [None] * n
    last = None
    for i in range(n - 1, -1, -1):
        if lines[i].strip():
            last = lines[i].strip()
        next_nonempty[i] = last
    for i, line in enumerate(lines):
        stripped = line.strip()
        ln = len(line) + 1  # +1 含换行
        if not stripped:
            pos += ln
            continue
        level = None
        if re.match(r"^#{1,6}\s+", stripped):
            level = len(re.match(r"^(#{1,6})", stripped).group(1))
        elif re.match(r"^第[一二三四五六七八九十百零\d]+\s*[章篇部部分]", stripped):
            level = 1
        elif re.match(r"^第[一二三四五六七八九十百零\d]+\s*[节课]", stripped):
            level = 2
        elif re.match(r"^第[一二三四五六七八九十百零\d]+\s*[条款项]", stripped):
            level = 3
        elif re.match(r"^Chapter\s+\d+", stripped, re.I):
            level = 1
        elif re.match(r"^Section\s+\d+", stripped, re.I):
            level = 2
        elif re.match(r"^Article\s+\d+", stripped, re.I):
            level = 3
        elif re.match(r"^\d+\.\d+\.\d+", stripped):
            level = 3
        elif re.match(r"^\d+\.\d+", stripped):
            level = 2
        elif re.match(r"^\d+\.", stripped):
            level = 1
        else:
            # 启发式：短、不以中英文句号结尾、后接长段落 -> 视作小标题
            nxt = next_nonempty[i] or ""
            if (len(stripped) <= 36 and not stripped[-1] in "。，.!?；："
                    and len(nxt) > 60):
                level = 3
        if level:
            headings.append({"level": level, "title": stripped, "char_start": pos})
        pos += ln
    return headings


def _locate_headings(structure, text):
    """给所有标题补上在正文中的字符偏移（docx/pptx 解析出的结构只有 level+title）。"""
    pos = 0
    out = []
    for h in structure:
        title = h.get("title", "")
        idx = text.find(title, pos) if title else -1
        if idx == -1:
            idx = pos
        out.append({"level": h.get("level", 3), "title": title, "char_start": idx})
        pos = idx + max(len(title), 1)
    out.sort(key=lambda x: x["char_start"])
    return out


def _detect_clauses(text):
    """条款检测：第X条 / Article X / Clause X / 多级编号 / （数字）编号。"""
    clauses = []
    patterns = [
        r"第[一二三四五六七八九十百零\d]+条",
        r"Article\s+\d+",
        r"Clause\s+\d+",
        r"^\(\d+\)",
        r"^（\d+）",
    ]
    for m in re.finditer("|".join(f"(?:{p})" for p in patterns), text, re.M | re.I):
        start = max(0, text.rfind("\n", 0, m.start()) + 1)
        end = text.find("\n", m.end())
        end = len(text) if end == -1 else end
        clauses.append({
            "label": m.group(0).strip(),
            "text": text[start:end].strip()[:300],
            "char_start": start,
        })
    return clauses


# ============================================================
# 结构化分块（百万字级不崩）
# ============================================================

def _build_chunks(text, structure):
    """按标题边界切块；超大块再按段落切分并加重叠。每块带标题面包屑。"""
    chunks = []
    if structure:
        bounds = [(h["char_start"], h["level"], h["title"]) for h in structure]
        bounds.append((len(text), 99, ""))  # 末尾哨兵
        # 层级栈维护面包屑：O(标题数)，避免对每块全量扫描 structure
        stack = []  # 元素为 (level, title)
        for idx in range(len(bounds) - 1):
            s, lvl, title = bounds[idx]
            e = bounds[idx + 1][0]
            while stack and stack[-1][0] >= lvl:
                stack.pop()
            stack.append((lvl, title))
            breadcrumb = [t for _, t in stack]
            sec = text[s:e].strip()
            chunks.extend(_split_long(sec, breadcrumb))
    else:
        # 无结构：按空行分段，合并小块
        blocks = [b.strip() for b in re.split(r"\n\s*\n", text) if b.strip()]
        buf, cur = [], ""
        for b in blocks:
            if len(cur) + len(b) > CHUNK_SIZE and cur:
                chunks.extend(_split_long(cur, []))
                cur = ""
            cur = (cur + "\n\n" + b) if cur else b
        if cur:
            chunks.extend(_split_long(cur, []))
    # 切块完成；向量在「检索时」用本地 _doc_vec 现算（不入库），避免大规模文档存巨量向量
    out = []
    for i, c in enumerate(chunks):
        if not c["text"].strip():
            continue
        c["i"] = i
        out.append(c)
    return out


def _split_long(sec, breadcrumb):
    """把超长文本按 CHUNK_SIZE 切分并带重叠。"""
    if len(sec) <= CHUNK_SIZE:
        return [{"text": sec, "headings": breadcrumb}]
    res = []
    start = 0
    while start < len(sec):
        end = start + CHUNK_SIZE
        piece = sec[start:end]
        res.append({"text": piece, "headings": breadcrumb})
        start = end - CHUNK_OVERLAP
        if start >= len(sec):
            break
    return res


# ============================================================
# 入库（上传入口）
# ============================================================

def ingest_file(path):
    """解析已保存到 uploads/ 的文件，并入文档库。返回该文档（含摘要）。"""
    _ensure_dirs()
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    dtype = EXT_MAP.get(ext, "txt")
    name = os.path.basename(path)
    size = os.path.getsize(path)

    # 按类型分发解析器
    if dtype == "pdf":
        text, tables, structure, meta = _parse_pdf(path)
    elif dtype == "docx":
        text, tables, structure, meta = _parse_docx(path)
    elif dtype == "xlsx":
        text, tables, structure, meta = _parse_xlsx(path)
    elif dtype == "pptx":
        text, tables, structure, meta = _parse_pptx(path)
    elif dtype == "csv":
        text, tables, structure, meta = _parse_csv(path)
    elif dtype == "html":
        text, tables, structure, meta = _parse_html(path, is_path=True)
    elif dtype == "md":
        text, tables, structure, meta = _parse_plain(path, as_markdown=True)
    elif dtype == "json":
        text, tables, structure, meta = _parse_plain(path, as_markdown=False)
    elif dtype == "image":
        text, tables, structure, meta = _parse_image(path)
    else:  # txt / text
        text, tables, structure, meta = _parse_plain(path, as_markdown=False)

    text = (text or "").strip()
    # 结构/条款若解析器没给（pdf/csv/html 等），用通用检测补全
    if not structure:
        structure = _detect_headings(text)
    structure = _locate_headings(structure, text)
    clauses = _detect_clauses(text)
    chunks = _build_chunks(text, structure)

    doc = {
        "id": uuid.uuid4().hex[:12],
        "name": name,
        "type": dtype,
        "size": size,
        "chars": len(text),
        "uploaded_at": datetime.datetime.now().isoformat(timespec="seconds"),
        "path": os.path.abspath(path),
        "full_text": text,
        "structure": structure,
        "tables": tables,
        "clauses": clauses,
        "chunks": chunks,
        "notes": meta.get("notes", ""),
        "pages": meta.get("pages"),
        "ocr": meta.get("ocr", False),
    }
    docs = _load_docs()
    docs[doc["id"]] = doc
    _save_docs(docs)
    return _summarize(doc)


def ingest_url(url):
    """下载网页 -> 存 uploads/ -> 走 HTML 解析。"""
    import httpx
    _ensure_dirs()
    try:
        r = httpx.get(url, timeout=20, follow_redirects=True,
                      headers={"User-Agent": "mini-agent/1.0"})
        r.raise_for_status()
        fname = re.sub(r"[^\w.-]", "_", url.split("?")[0].split("/")[-1] or "web")
        if not fname.lower().endswith((".html", ".htm")):
            fname += ".html"
        path = os.path.join(UPLOAD_DIR, fname)
        with open(path, "wb") as f:
            f.write(r.content)
    except Exception as e:
        return {"error": f"下载网页失败：{e}"}
    doc = ingest_file(path)
    doc["source_url"] = url
    return doc


# ============================================================
# 对外查询 API
# ============================================================

def _summarize(doc):
    return {
        "id": doc["id"], "name": doc["name"], "type": doc["type"],
        "chars": doc["chars"], "pages": doc.get("pages"),
        "tables": len(doc["tables"]), "chunks": len(doc["chunks"]),
        "clauses": len(doc["clauses"]), "headings": len(doc["structure"]),
        "notes": doc.get("notes", ""), "ocr": doc.get("ocr", False),
    }


def list_documents():
    docs = _load_docs()
    if not docs:
        return "还没有上传任何文档。上传后会自动解析、切块、建索引。"
    lines = ["已上传文档："]
    for d in docs.values():
        s = _summarize(d)
        note = f"（提示：{s['notes']}）" if s["notes"] else ""
        ocr = " [OCR]" if s["ocr"] else ""
        lines.append(
            f"- 【{s['type']}{ocr}】{s['name']}  id={s['id']}\n"
            f"  字数 {s['chars']} · 表格 {s['tables']} · 分块 {s['chunks']} · "
            f"条款 {s['clauses']} · 标题 {s['headings']}{note}"
        )
    return "\n".join(lines)


def list_document_summaries():
    """返回结构化文档摘要列表（供前端渲染）。"""
    return [_summarize(d) for d in _load_docs().values()]


def get_document(doc_id):
    docs = _load_docs()
    return docs.get(doc_id)


def read_document(doc_id, max_chars=4000):
    doc = get_document(doc_id)
    if not doc:
        return f"找不到文档 id={doc_id}"
    txt = doc["full_text"][:max_chars]
    more = "" if doc["chars"] <= max_chars else f"\n…（共 {doc['chars']} 字，已截断，可用 search_documents 精准定位）"
    return f"文档《{doc['name']}》（{doc['chars']} 字，{len(doc['chunks'])} 块）：\n{txt}{more}"


def extract_tables(doc_id):
    doc = get_document(doc_id)
    if not doc:
        return f"找不到文档 id={doc_id}"
    if not doc["tables"]:
        return f"《{doc['name']}》中未检测到表格。"
    out = [f"《{doc['name']}》共 {len(doc['tables'])} 张表："]
    for i, t in enumerate(doc["tables"], 1):
        cap = t.get("caption") or f"表 {i}"
        out.append(f"\n### {cap}")
        out.append(" | ".join(str(h) for h in t["headers"]))
        out.append("-" * 30)
        for row in t["rows"][:30]:
            out.append(" | ".join(str(c) for c in row))
    if len(doc["tables"]) > 0:
        out.append(f"\n（最多展示每张表前 30 行）")
    return "\n".join(out)


def extract_clauses(doc_id, keyword=None):
    doc = get_document(doc_id)
    if not doc:
        return f"找不到文档 id={doc_id}"
    clauses = doc["clauses"]
    if keyword:
        kw = keyword.lower()
        clauses = [c for c in clauses if kw in c["text"].lower()]
    if not clauses:
        return f"《{doc['name']}》中未检测到条款" + (f"（关键词：{keyword}）" if keyword else "")
    out = [f"《{doc['name']}》检测到 {len(clauses)} 条相关条款："]
    for c in clauses[:50]:
        out.append(f"\n【{c['label']}】{c['text']}")
    return "\n".join(out)


def search_documents(query, top_k=5):
    """跨所有文档的分块做关键词检索，返回最相关片段。"""
    docs = _load_docs()
    if not docs:
        return "还没有文档，先上传再搜索。"
    q = _doc_vec(query)
    scored = []
    for doc in docs.values():
        for ch in doc["chunks"]:
            s = _cosine(q, _doc_vec(ch["text"]))
            if s > 0:
                scored.append((s, doc["name"], doc["id"], ch))
    scored.sort(key=lambda x: x[0], reverse=True)
    if not scored:
        return "没有找到相关内容。"
    # 关键词向量的相似度分布偏低，用较松的阈值
    thr = 0.04
    top = [x for x in scored if x[0] >= thr][:top_k]
    if not top:
        top = scored[:top_k]
    lines = [f"跨文档检索「{query}」命中 {len(top)} 处："]
    for s, name, did, ch in top:
        bc = " > ".join(ch.get("headings", [])[-3:]) if ch.get("headings") else ""
        snippet = ch["text"][:240].replace("\n", " ")
        lines.append(
            f"\n· 《{name}》(id={did}) 相关度 {s:.2f}"
            + (f"  [{bc}]" if bc else "")
            + f"\n  {snippet}"
        )
    return "\n".join(lines)


def compare_documents(a_id, b_id, topic=None):
    """跨文档关联比对 + 交叉验证（启发式）。
    返回：最相似片段配对（重叠/一致候选）+ 数值差异扫描（建议人工/LLM 复核）。"""
    docs = _load_docs()
    da, db = docs.get(a_id), docs.get(b_id)
    if not da or not db:
        return f"比对需要两个有效文档 id，当前：a={a_id}, b={b_id}"
    na, nb = da["name"], db["name"]

    # 1) 相似配对：B 中每段找 A 中最佳匹配
    pairs = []
    pool = [(_doc_vec(ch["text"]), ch["text"]) for ch in db["chunks"]]
    for ch in da["chunks"]:
        best, best_s = None, -1
        vec_a = _doc_vec(ch["text"])
        for vec, txt in pool:
            s = _cosine(vec_a, vec)
            if s > best_s:
                best_s, best = s, txt
        if best is not None:
            pairs.append((best_s, ch["text"], best))
    pairs.sort(key=lambda x: x[0], reverse=True)
    out = [f"跨文档比对：《{na}》 ↔ 《{nb}》",
           f"（比对 {len(da['chunks'])}×{len(db['chunks'])} 个分块，以下为最相似片段：）"]

    # 若给了主题，先用检索聚焦
    focus = []
    if topic:
        q = _doc_vec(topic)
        for doc, label in ((da, na), (db, nb)):
            for ch in doc["chunks"]:
                sc = _cosine(q, _doc_vec(ch["text"]))
                if sc >= 0.04:
                    focus.append((label, ch["text"][:200]))
        if focus:
            out.append(f"\n【主题「{topic}」相关片段】")
            for label, txt in focus[:8]:
                out.append(f"- 《{label}》：{txt}")

    out.append("\n【最相似片段配对（一致性/重叠候选）】")
    for s, ta, tb in pairs[:5]:
        out.append(f"\n 相似度 {s:.2f}")
        out.append(f"  A: {ta[:180].replace(chr(10),' ')}")
        out.append(f"  B: {tb[:180].replace(chr(10),' ')}")

    # 2) 数值差异扫描（轻量启发式，仅作提示）
    diffs = _scan_numeric_divergence(da["full_text"], db["full_text"])
    if diffs:
        out.append("\n【数值差异提示（建议人工/LLM 复核是否冲突）】")
        out.extend(f"  - {d}" for d in diffs[:10])

    out.append("\n（注：以上为启发式比对；配置真实大模型后，compare_documents 可由 LLM 做深层语义冲突判定。）")
    return "\n".join(out)


def _scan_numeric_divergence(ta, tb):
    """抽取『数量+单位』，若同一单位在两边数值差异大，则标记。仅作粗筛。"""
    pat = re.compile(r"(\d[\d,\.]*)\s*(元|万元|亿美元|美元|%|％|kg|千克|吨|个|条|人|项|年|月|天|平方米|平米|m²)")
    def collect(t):
        d = {}
        for num, unit in pat.findall(t):
            v = float(num.replace(",", ""))
            d.setdefault(unit, []).append(v)
        return d
    a, b = collect(ta), collect(tb)
    diffs = []
    for unit in set(a) & set(b):
        av, bv = a[unit], b[unit]
        if max(av) / max(bv) > 1.5 or min(av) / max(bv) < 0.67:
            diffs.append(f"单位「{unit}」：A 取值 {sorted(av)} vs B 取值 {sorted(bv)}")
    return diffs
